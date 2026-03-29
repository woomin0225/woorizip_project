from __future__ import annotations

import base64
import binascii
import io
import re
import zipfile
import zlib
from typing import Any
from xml.etree import ElementTree as ET


TEXT_MIME_TYPES = {
    'text/plain',
    'text/markdown',
    'text/csv',
    'application/json',
}

PDF_MIME_TYPES = {'application/pdf'}
DOCX_MIME_TYPES = {
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
}
XLSX_MIME_TYPES = {
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
}
PPTX_MIME_TYPES = {
    'application/vnd.openxmlformats-officedocument.presentationml.presentation',
}
HWPX_MIME_TYPES = {
    'application/hwpx',
    'application/x-hwpx',
    'application/haansofthwpx',
}
HWP_MIME_TYPES = {
    'application/hwp',
    'application/x-hwp',
    'application/haansofthwp',
}


class DocumentParseError(Exception):
    pass


class DocumentParser:
    def parse_attachment(self, attachment: dict[str, Any]) -> dict[str, Any]:
        filename = attachment.get('filename') or 'unknown'
        mime_type = (attachment.get('mime_type') or '').strip().lower()
        raw_text = attachment.get('text')
        base64_data = attachment.get('file_base64')

        if raw_text:
            text = self._normalize_text(str(raw_text))
            return {
                'filename': filename,
                'mime_type': mime_type or 'text/plain',
                'text': text,
                'chars': len(text),
                'source': 'text',
            }

        if not base64_data:
            raise DocumentParseError('text 또는 file_base64 중 하나는 필요합니다.')

        try:
            file_bytes = base64.b64decode(base64_data, validate=True)
        except (binascii.Error, ValueError) as exc:
            raise DocumentParseError('base64 디코딩에 실패했습니다.') from exc

        text = self._extract_text(file_bytes, mime_type=mime_type, filename=filename)
        return {
            'filename': filename,
            'mime_type': mime_type or self._infer_mime_from_filename(filename),
            'text': text,
            'chars': len(text),
            'source': 'file',
        }

    def _extract_text(self, file_bytes: bytes, mime_type: str, filename: str) -> str:
        effective_mime = mime_type or self._infer_mime_from_filename(filename)
        lowered = filename.lower()

        if effective_mime in TEXT_MIME_TYPES:
            return self._normalize_text(file_bytes.decode('utf-8', errors='ignore'))

        if effective_mime in PDF_MIME_TYPES or lowered.endswith('.pdf'):
            return self._extract_pdf_text(file_bytes)

        if effective_mime in DOCX_MIME_TYPES or lowered.endswith('.docx'):
            return self._extract_docx_text(file_bytes)

        if effective_mime in XLSX_MIME_TYPES or lowered.endswith('.xlsx'):
            return self._extract_xlsx_text(file_bytes)

        if effective_mime in PPTX_MIME_TYPES or lowered.endswith('.pptx'):
            return self._extract_pptx_text(file_bytes)

        if effective_mime in HWPX_MIME_TYPES or lowered.endswith('.hwpx'):
            return self._extract_hwpx_text(file_bytes)

        if effective_mime in HWP_MIME_TYPES or lowered.endswith('.hwp'):
            return self._extract_hwp_text(file_bytes)

        raise DocumentParseError(f'지원하지 않는 첨부 형식입니다: {mime_type or filename}')

    def _extract_pdf_text(self, file_bytes: bytes) -> str:
        try:
            import fitz  # PyMuPDF
        except ImportError as exc:
            raise DocumentParseError('PyMuPDF가 설치되어 있지 않습니다.') from exc

        parts: list[str] = []
        try:
            with fitz.open(stream=file_bytes, filetype='pdf') as doc:
                for page in doc:
                    parts.append(page.get_text('text'))
        except Exception as exc:  # pragma: no cover
            raise DocumentParseError('PDF 텍스트 추출에 실패했습니다.') from exc
        return self._normalize_text('\n'.join(parts))

    def _extract_docx_text(self, file_bytes: bytes) -> str:
        try:
            from docx import Document
        except ImportError as exc:
            raise DocumentParseError('python-docx가 설치되어 있지 않습니다.') from exc

        try:
            doc = Document(io.BytesIO(file_bytes))
            paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
        except Exception as exc:  # pragma: no cover
            raise DocumentParseError('DOCX 텍스트 추출에 실패했습니다.') from exc
        return self._normalize_text('\n'.join(paragraphs))

    def _extract_xlsx_text(self, file_bytes: bytes) -> str:
        try:
            from openpyxl import load_workbook
        except ImportError as exc:
            raise DocumentParseError('openpyxl이 설치되어 있지 않습니다.') from exc

        sections: list[str] = []

        try:
            workbook = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
            for sheet in workbook.worksheets:
                rows: list[str] = []
                for row in sheet.iter_rows(values_only=True):
                    values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                    if values:
                        rows.append('\t'.join(values))

                if rows:
                    sections.append(f"[시트] {sheet.title}\n" + '\n'.join(rows))
        except Exception as exc:  # pragma: no cover
            raise DocumentParseError('XLSX 텍스트 추출에 실패했습니다.') from exc

        if not sections:
            raise DocumentParseError('XLSX 파일에서 추출할 텍스트를 찾지 못했습니다.')

        return self._normalize_text('\n\n'.join(sections))

    def _extract_pptx_text(self, file_bytes: bytes) -> str:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise DocumentParseError('python-pptx가 설치되어 있지 않습니다.') from exc

        sections: list[str] = []

        try:
            presentation = Presentation(io.BytesIO(file_bytes))
            for index, slide in enumerate(presentation.slides, start=1):
                texts: list[str] = []

                for shape in slide.shapes:
                    if getattr(shape, 'has_text_frame', False):
                        text = getattr(shape, 'text', '')
                        if text and text.strip():
                            texts.append(text.strip())
                        continue

                    if getattr(shape, 'has_table', False):
                        for row in shape.table.rows:
                            row_values = [
                                cell.text.strip()
                                for cell in row.cells
                                if cell.text and cell.text.strip()
                            ]
                            if row_values:
                                texts.append('\t'.join(row_values))

                if texts:
                    sections.append(f"[슬라이드 {index}]\n" + '\n'.join(texts))
        except Exception as exc:  # pragma: no cover
            raise DocumentParseError('PPTX 텍스트 추출에 실패했습니다.') from exc

        if not sections:
            raise DocumentParseError('PPTX 파일에서 추출할 텍스트를 찾지 못했습니다.')

        return self._normalize_text('\n\n'.join(sections))

    def _extract_hwpx_text(self, file_bytes: bytes) -> str:
        paragraph_texts: list[str] = []

        try:
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as archive:
                xml_names = sorted(
                    name
                    for name in archive.namelist()
                    if name.lower().endswith('.xml') and '/section' in name.lower()
                )

                if not xml_names:
                    xml_names = sorted(
                        name
                        for name in archive.namelist()
                        if name.lower().endswith('.xml') and 'contents/' in name.lower()
                    )

                for xml_name in xml_names:
                    content = archive.read(xml_name)
                    root = ET.fromstring(content)

                    for paragraph in root.iter():
                        if self._local_name(paragraph.tag) != 'p':
                            continue

                        parts: list[str] = []
                        for node in paragraph.iter():
                            if self._local_name(node.tag) != 't':
                                continue
                            value = ''.join(node.itertext()).strip()
                            if value:
                                parts.append(value)

                        if parts:
                            paragraph_texts.append(''.join(parts))
        except zipfile.BadZipFile as exc:
            raise DocumentParseError('HWPX 파일 형식이 올바르지 않습니다.') from exc
        except Exception as exc:  # pragma: no cover
            raise DocumentParseError('HWPX 텍스트 추출에 실패했습니다.') from exc

        if not paragraph_texts:
            raise DocumentParseError('HWPX 파일에서 추출할 텍스트를 찾지 못했습니다.')

        return self._normalize_text('\n'.join(paragraph_texts))

    def _extract_hwp_text(self, file_bytes: bytes) -> str:
        try:
            import olefile
        except ImportError as exc:
            raise DocumentParseError('olefile이 설치되어 있지 않습니다.') from exc

        ole = None
        try:
            ole = olefile.OleFileIO(io.BytesIO(file_bytes))
            preview_text = self._extract_hwp_preview_text(ole)
            if preview_text:
                return preview_text

            body_text = self._extract_hwp_body_text(ole)
            if body_text:
                return body_text
        except DocumentParseError:
            raise
        except Exception as exc:  # pragma: no cover
            raise DocumentParseError('HWP 텍스트 추출에 실패했습니다.') from exc
        finally:
            if ole is not None:
                ole.close()

        raise DocumentParseError('HWP 파일에서 추출할 텍스트를 찾지 못했습니다.')

    def _extract_hwp_preview_text(self, ole) -> str:
        if not ole.exists(['PrvText']):
            return ''

        preview_bytes = ole.openstream(['PrvText']).read()
        for encoding in ('utf-16le', 'utf-16', 'utf-8'):
            try:
                text = self._normalize_text(preview_bytes.decode(encoding))
            except UnicodeDecodeError:
                continue
            if text:
                return text
        return ''

    def _extract_hwp_body_text(self, ole) -> str:
        compressed = self._is_hwp_body_compressed(ole)
        sections = [
            entry
            for entry in ole.listdir(streams=True, storages=False)
            if len(entry) == 2 and entry[0] == 'BodyText' and entry[1].startswith('Section')
        ]
        sections.sort(key=self._hwp_section_sort_key)

        texts: list[str] = []
        for section in sections:
            section_bytes = ole.openstream(section).read()
            if compressed:
                try:
                    section_bytes = zlib.decompress(section_bytes, -15)
                except zlib.error:
                    continue

            section_text = self._extract_hwp_section_text(section_bytes)
            if section_text:
                texts.append(section_text)

        return self._normalize_text('\n'.join(texts))

    def _is_hwp_body_compressed(self, ole) -> bool:
        if not ole.exists(['FileHeader']):
            return False

        header = ole.openstream(['FileHeader']).read()
        if len(header) < 40:
            return False

        properties = int.from_bytes(header[36:40], 'little', signed=False)
        return bool(properties & 0x01)

    def _hwp_section_sort_key(self, entry: list[str]) -> int:
        match = re.search(r'(\d+)$', entry[1])
        return int(match.group(1)) if match else 0

    def _extract_hwp_section_text(self, section_bytes: bytes) -> str:
        runs: list[str] = []
        current: list[str] = []

        for index in range(0, len(section_bytes) - 1, 2):
            code_point = section_bytes[index] | (section_bytes[index + 1] << 8)
            char = chr(code_point)

            if self._is_reasonable_hwp_char(char):
                current.append(char)
                continue

            self._flush_hwp_text_run(current, runs)

        self._flush_hwp_text_run(current, runs)
        return self._normalize_text('\n'.join(runs))

    def _flush_hwp_text_run(self, current: list[str], runs: list[str]) -> None:
        if not current:
            return

        text = ''.join(current)
        current.clear()

        normalized = self._normalize_text(text)
        if not normalized:
            return

        if not any(self._is_meaningful_hwp_char(char) for char in normalized):
            return

        if len(normalized.replace(' ', '')) < 2:
            return

        runs.append(normalized)

    def _is_reasonable_hwp_char(self, char: str) -> bool:
        if char in {'\n', '\r', '\t', ' '}:
            return True

        code = ord(char)
        if 0xAC00 <= code <= 0xD7A3:
            return True
        if 0x3131 <= code <= 0x318E:
            return True
        if 0x1100 <= code <= 0x11FF:
            return True
        if char.isascii() and (char.isalnum() or char in '.,!?()[]{}<>+-_/:%&@#*\'"~'):
            return True

        return False

    def _is_meaningful_hwp_char(self, char: str) -> bool:
        code = ord(char)
        if 0xAC00 <= code <= 0xD7A3:
            return True
        if 0x3131 <= code <= 0x318E:
            return True
        return char.isalnum()

    def _normalize_text(self, text: str) -> str:
        lines = [line.strip() for line in text.replace('\r', '\n').split('\n')]
        compact = [line for line in lines if line]
        return '\n'.join(compact)

    def _local_name(self, tag: str) -> str:
        if not tag:
            return ''
        if '}' in tag:
            return tag.rsplit('}', 1)[-1]
        return tag

    def _infer_mime_from_filename(self, filename: str) -> str:
        lowered = filename.lower()
        if lowered.endswith('.pdf'):
            return 'application/pdf'
        if lowered.endswith('.docx'):
            return 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        if lowered.endswith('.xlsx'):
            return 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        if lowered.endswith('.pptx'):
            return 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        if lowered.endswith('.hwpx'):
            return 'application/haansofthwpx'
        if lowered.endswith('.hwp'):
            return 'application/x-hwp'
        if lowered.endswith('.txt'):
            return 'text/plain'
        if lowered.endswith('.md'):
            return 'text/markdown'
        return 'application/octet-stream'
